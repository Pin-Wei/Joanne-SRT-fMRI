#!/usr/bin/env python

import os
import re
import glob
import argparse
import numpy as np
import pandas as pd
from scipy.io import loadmat

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
    from pptx.enum.shapes import PP_PLACEHOLDER
    from pptx.oxml.xmlchemy import OxmlElement
    from pptx.oxml.ns import qn
except ImportError:
    print("The 'python-pptx' library is required to run this script." + 
          "\nPlease install it using 'conda/mamba install -c conda-forge python-pptx'.")
    exit(1)

try:
    from PIL import Image
except ImportError:
    print("The 'Pillow' library is required to run this script." + 
          "\nPlease install it using 'conda/mamba install -c conda-forge pillow'.")
    exit(1)

class Config:
    def __init__(self, args):
        self.preset = {1: "FNC", 2: "SPC", 3: "TFCE"}[args.preset]
        self.preset_full = {
            1: "Functional Network Connectivity", # (Jafri et al., 2008)
            2: "Spatial Pairwise Clustering", # (Zalesky et al., 2012)
            3: "Threshold-Free Cluster Enhancement" # (Smith & Nichols, 2007)
        }[args.preset]
        self.title = "ROI-to-ROI gPPI results"
        self.subtitle = (
            # f"Project: {args.project}\n" + 
            f"ROI prefix: {args.roi_prefix}\n" +
            f"Inference preset: {self.preset} ({self.preset_full})"
        )
        self.setup_paths(args)

    def setup_paths(self, args):
        self.source_dir = os.path.dirname(os.path.abspath(__file__))
        self.mat_file = os.path.join(self.source_dir, "..", "conn_out", f"{args.project}.mat")        
        self.qa_dir = os.path.join(self.source_dir, "..", "conn_out", args.project, "results", "qa")
        dv_files = glob.glob(os.path.join(self.qa_dir, "QA_GUIrequest_DataValidity_*", "DataValidityScore.mat"))
        self.dv_file = dv_files[0] if len(dv_files) > 0 else None
        dq_files = glob.glob(os.path.join(self.qa_dir, "QA_GUIrequest_DataQuality_*", "DataQualityScore.mat"))
        self.dq_file = dq_files[0] if len(dq_files) > 0 else None
        ds_files = glob.glob(os.path.join(self.qa_dir, "QA_GUIrequest_DataSensitivity_*", "DataSensitivityScore.mat"))
        self.ds_file = ds_files[0] if len(ds_files) > 0 else None        
        self.report_dir = os.path.join(self.source_dir, "..", "reports", args.project, f"ROI-{args.roi_prefix}_{self.preset}")
        self.report_path = os.path.join(self.report_dir, "report.pptx")

        while os.path.exists(self.report_path) and not args.overwrite:
            resp = input(f"'{self.report_path}' already exists. Do you want to overwrite it? (y/n): ")
            if resp.lower() == "y":
                print("Overwriting existing report ...")
                break
            else:
                resp = input("Do you want to specify a different filename? (y/n): ")
                if resp.lower() == "y":
                    new_name = input("Enter the new filename (without extension): ")
                    self.report_path = os.path.join(self.report_dir, f"{new_name}.pptx")
                else:
                    print("Exiting without overwriting ...")
                    exit(0)

class Constants:
    def __init__(self):
        self.slide_width = Inches(13.33)
        self.slide_height = Inches(7.5)
        self.setup_colors()
        self.setup_master()
        self.setup_table()

        self.img_height = Inches(5.5)
        self.foot_font_size = Pt(18)
        self.foot_pos = [
            Inches(0.5), 
            self.slide_height - Inches(0.6),
            Inches(5),
            Inches(0.3)
        ]
        self.slide_num_font_size = Pt(16)
        self.slide_num_pos = [
            Inches(12.5), 
            self.slide_height - Inches(0.5), 
            Inches(0.5),
            Inches(0.3)
        ]

    def setup_colors(self):
        self.bullet_color = RGBColor(192, 79, 21) # Burnt Orange
        self.str_title_color = RGBColor(0, 112, 192) # Royal Navy Blue
        self.swi_title_color = RGBColor(0, 176, 80) # Muslim Green
        self.foot_color = RGBColor(0, 0, 0)

    def setup_master(self):
        self.master_font_name = "Calibri"
        self.master_l0_title_font_size = Pt(60) # layout 0 == 'Title Slide'
        self.master_l0_title_pos = [
            Inches(1.67), # left
            Inches(1.23), # top
            Inches(10),   # width
            Inches(2.61)  # height
        ]
        self.master_l0_subtitle_font_size = Pt(24)
        self.master_l0_subtitle_pos = [
            Inches(1.67), 
            Inches(3.94),
            Inches(10), 
            Inches(1.81)
        ]
        self.master_l1_title_font_size = Pt(36) # layout 1 == 'Title and Content'
        self.master_l1_title_pos = [
            Inches(0.56), 
            Inches(0.4), 
            Inches(12.22), 
            Inches(0.66)
        ]
        self.master_l1_content_font_size = Pt(24)
        self.master_l1_content_pos = [
            Inches(0.56), 
            Inches(1.24), 
            Inches(12.22), 
            Inches(5.52)
        ]

    def setup_table(self):
        self.tbl_row_height = Inches(0.3)
        self.tbl_col_widths = [
            Inches(0.5), 
            Inches(3), 
            Inches(3), 
            Inches(1.75), 
            Inches(1), 
            Inches(1), 
            Inches(1)
        ]
        self.tbl_width = sum(self.tbl_col_widths)
        self.tbl_left = (self.slide_width - self.tbl_width) // 2 # center the table horizontally

def parse_args():
    parser = argparse.ArgumentParser()
    parser.add_argument("-p", "--project", type=str, required=True, 
                        help="Name of the CONN project.")
    parser.add_argument("-r", "--roi_prefix", type=str, default="networks", 
                        help="Prefix of the ROI names used to generate the report.")
    parser.add_argument("-i", "--preset", type=int, choices=[1, 2, 3], default=3,
                        help="Inference preset used to generate the report. (1: FNC, 2: SPC, 3: TFCE)")
    parser.add_argument("-o", "--overwrite", action="store_true", 
                        help="Overwrite existing files.")
    return parser.parse_args()

def setup_and_get_master_layout(prs, layout_name, const):
    '''
    Set up the master layout (placeholder positions and font sizes), 
    and return the layout object.
    '''
    def _get_or_add(parent, tag):
        child = parent.find(qn(tag))
        if child is None:
            child = OxmlElement(tag)
            parent.append(child)
        return child

    def _set_ph(shape, pos, font_size, align, anchor, 
                bold=False, italic=False, underline=False, is_slide_num=False):
        shape.left, shape.top, shape.width, shape.height = pos

        txBody = shape._element.txBody
        bodyPr = _get_or_add(txBody, "a:bodyPr")
        bodyPr.set("anchor", {"top": "t", "center": "ctr", "bottom": "b"}[anchor]) # vertical alignment

        lstStyle = _get_or_add(txBody, "a:lstStyle")
        lvl1pPr = _get_or_add(lstStyle, "a:lvl1pPr")
        lvl1pPr.set("algn", {"left": "l", "center": "ctr", "right": "r"}[align]) # horizontal alignment

        defRPr = _get_or_add(lvl1pPr, "a:defRPr")
        defRPr.set("sz", str(int(font_size.pt * 100))) # pt -> twips

        if bold:
            defRPr.set("b", "1")
        else:
            defRPr.attrib.pop("b", None)

        if italic:
            defRPr.set("i", "1")
        else:
            defRPr.attrib.pop("i", None)

        if underline:
            defRPr.set("u", "sng")
        else:
            defRPr.attrib.pop("u", None)

        latin = _get_or_add(defRPr, "a:latin")
        latin.set("typeface", const.master_font_name)

        if not is_slide_num:
            shape.text_frame.clear() # avoid leaving sample text in the layout placeholder

    master = prs.slide_master # since most presentations have only a single slide master

    for layout in master.slide_layouts:
        if layout.name != layout_name: # find the layout by name
            continue
        
        for shape in layout.shapes: 
            # print(f"Processing {shape.name} ...")

            if layout.name == "Title Slide":
                if "Title" in shape.name:
                    _set_ph(
                        shape, 
                        pos=const.master_l0_title_pos, 
                        font_size=const.master_l0_title_font_size, 
                        align="center", 
                        anchor="bottom", 
                        bold=True
                    )

                elif "Subtitle" in shape.name:
                    _set_ph(
                        shape, 
                        pos=const.master_l0_subtitle_pos, 
                        font_size=const.master_l0_subtitle_font_size, 
                        align="center", 
                        anchor="top"
                    )

            elif layout.name in ["Title and Content", "Title Only"]:
                if "Title" in shape.name:
                    _set_ph(
                        shape, 
                        pos=const.master_l1_title_pos, 
                        font_size=const.master_l1_title_font_size, 
                        align="left", 
                        anchor="bottom"
                    )

                elif "Content" in shape.name:
                    _set_ph(
                        shape, 
                        pos=const.master_l1_content_pos, 
                        font_size=const.master_l1_content_font_size, 
                        align="left", 
                        anchor="top"
                    )

            if "Slide Number" in shape.name:
                _set_ph(
                    shape, 
                    pos=const.slide_num_pos, 
                    font_size=const.slide_num_font_size, 
                    align="right", 
                    anchor="bottom", 
                    italic=True, 
                    is_slide_num=True
                )
            
        return layout
    
    raise ValueError(f"Layout '{layout_name}' not found.")

def add_bullet_points(tf, lines, const):
    '''
    Add bullet points to the text frame.
    '''
    for i, line in enumerate(lines):
        if ": " in line:
            txt1, txt2 = line.split(": ")
            txt1 += ": "
        else:
            txt1, txt2 = None, line

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

        if txt1:
            r1 = p.add_run()
            r1.text = txt1
            r1.font.color.rgb = const.bullet_color

        r2 = p.add_run()
        r2.text = txt2
        r2.font.color.rgb = RGBColor(0, 0, 0)

def get_denoising_params(mat_file):
    '''
    Extract preprocessing and denoising parameters from the CONN .mat file.
    '''
    data = loadmat(mat_file, simplify_cells=True)
    confound_names = data["CONN_x"]["Preproc"]["confounds"]["names"]
    confound_powers = data["CONN_x"]["Preproc"]["confounds"]["power"]
    confound_dims = data["CONN_x"]["Preproc"]["confounds"]["dimensions"]
    mot_idx = np.where(confound_names == "realignment")[0][0]
    acomp_idx = np.where(confound_names == "aCompCor")[0][0]
    aroma_idx = np.where(confound_names == "aroma")[0][0]
    gmr_idx = np.where(confound_names == "Gray Matter")[0]

    return {
        "POLY_ORD": data["CONN_x"]["Preproc"]["detrending"],
        "BP_LO": data["CONN_x"]["Preproc"]["filter"][0],
        "BP_HI": data["CONN_x"]["Preproc"]["filter"][1],
        "SIMULT": data["CONN_x"]["Preproc"]["regbp"] == 0,
        "MOT24": confound_powers[mot_idx] == 2,
        "N_ACOMP": confound_dims[acomp_idx][0],
        "N_AROMA": confound_dims[aroma_idx][0],
        "RM_GMR": len(gmr_idx) > 0 
    }

def get_qa_scores(config):
    '''
    Extract QA scores (Data Validity, Quality, Sensitivity) from the CONN .mat files.
    '''
    scores = {}
    for score_name, fp in [
        ("Data Validity (DV) score", config.dv_file), 
        ("Data Quality (DQ) score", config.dq_file), 
        ("Data Sensitivity (DS) score", config.ds_file)
    ]:
        if fp:
            data = loadmat(fp, simplify_cells=True)
            fn = os.path.basename(fp).split(".")[0]
            score = data[fn]
            scores[score_name] = score[0] if isinstance(score, np.ndarray) else score
        else:
            scores[score_name] = None

    return scores

def add_contrast_title(tf, contrast, view_name, const):
    '''
    Add a formatted title for the contrast (and view name) to the given text frame.
    '''
    def _add_run_with_color(p, contrast, color):
        runs_str = re.findall(r"_r([\d]+)", contrast)
        runs = f'run-{", ".join([ r for r in runs_str[0] ])}' if runs_str else "all"
        r2 = p.add_run()
        r2.text = runs
        r2._r.get_or_add_rPr().set("baseline", "-25000") # subscripts
        r2.font.color.rgb = color
        r2.font.bold = True
        r2.font.italic = True
    
    def _add_minus_random(p):
        r3 = p.add_run()
        r3.text = " minus"
        r3.font.color.rgb = RGBColor(0, 0, 0)
        r4 = p.add_run()
        r4.text = " Random"
        r4.font.color.rgb = RGBColor(0, 0, 0)
        r4.font.bold = True

    def _add_view_name(p, view_name):
        r5 = p.add_run()
        r5.text = f" ({view_name})"
        r5.font.size = Pt(24)
        r5.font.italic = True
        r5.font.color.rgb = RGBColor(0, 0, 0)

    p = tf.paragraphs[0]
    r1 = p.add_run()

    if contrast.startswith("str"):
        if "fst" in contrast:
            r1.text = "Sequence-1"
        elif "snd" in contrast:
            r1.text = "Sequence-2"
        else:
            r1.text = "Structured"

        r1.font.bold = True
        r1.font.color.rgb = const.str_title_color
        _add_run_with_color(p, contrast, const.str_title_color)
        _add_minus_random(p)
        
    elif contrast.startswith("swi"):
        r1.text = "Switch"
        r1.font.bold = True
        r1.font.color.rgb = const.swi_title_color
        _add_run_with_color(p, contrast, const.swi_title_color)
        _add_minus_random(p)
    
    else: # undefined contrast type, just add the raw text
        r1.text = contrast

    if view_name:
        _add_view_name(p, view_name)
     
def add_stats_table(slide, stats_df, const):
    '''
    Add a formatted table of statistics to the given slide, based on the stats DataFrame.
    '''
    def _format_text(x):
        if pd.isna(x):
            return ""
        elif isinstance(x, float):
            return f"{x:.3f}"
        else:
            x = str(x)
            parts = re.findall(r"(.*)\s+\(-?[\d]+[\s]+-?[\d]+[\s]+-?[\d]+\)", x)
            if parts:
                return parts[0]
            else:
                return x

    n_rows, n_cols = stats_df.shape[0] + 1, stats_df.shape[1]
    assert n_cols == len(const.tbl_col_widths), f"Expected {len(const.tbl_col_widths)} columns in the stats CSV."
    
    tbl_height = const.tbl_row_height * n_rows
    tbl_top = (const.slide_num_pos[1] - const.master_l1_content_pos[1] - tbl_height) // 2 
    tbl_top = const.master_l1_content_pos[1] if tbl_top < const.master_l1_content_pos[1] else tbl_top

    table = slide.shapes.add_table(
        n_rows, n_cols, 
        const.tbl_left, tbl_top, const.tbl_width, tbl_height
    ).table

    # Set column widths and fill headers
    for col_idx, col_name in enumerate(stats_df.columns):
        table.columns[col_idx].width = const.tbl_col_widths[col_idx]

        if col_idx != 0: # skip the first column
            table.cell(0, col_idx).text = str(col_name)
            table.cell(0, col_idx).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Write body cells
    prev_cid = None

    for row_idx, (_, row) in enumerate(stats_df.iterrows()):
        ridx = row_idx + 1

        for col_idx, value in enumerate(row):
            cell = table.cell(ridx, col_idx)

            if col_idx == 0: # cluster ID row
                cid = int(value)

                if cid != prev_cid:
                    cell.text = str(cid)
                    p = cell.text_frame.paragraphs[0]
                    p.alignment = PP_ALIGN.CENTER
                    p.font.italic = True

                prev_cid = cid
            else:
                cell.text = _format_text(value)

def add_figure_and_footnote(slide, fig_path, footnote, const):
    '''
    Add a figure and a footnote to the given slide.
    '''
    with Image.open(fig_path) as img:
        img_width, img_height = img.size
    
    img_aspect = img_width / img_height
    img_height = const.img_height 
    img_width = const.img_height * img_aspect
    img_left = (const.slide_width - img_width) // 2
    img_top = const.master_l1_content_pos[1] 
    
    slide.shapes.add_picture(
        fig_path, 
        img_left, img_top, img_width, img_height
    )

    # Add footnote
    left, top, width, height = const.foot_pos
    foot = slide.shapes.add_textbox(left, top, width, height)
    tf = foot.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = footnote
    p.font.size = const.foot_font_size
    p.font.color.rgb = const.foot_color
    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

# def add_slide_number(prs, const):
    # for i, slide in enumerate(prs.slides):
    #     tb = slide.shapes.add_textbox(
    #         const.slide_num_pos[0], const.slide_num_pos[1], const.slide_num_pos[2], const.slide_num_pos[3]
    #     )
    #     tf = tb.text_frame
    #     p = tf.paragraphs[0]
    #     p.text = str(i + 1)
    #     p.font.size = const.slide_num_font_size
    #     p.font.italic = True
    #     p.alignment = PP_ALIGN.RIGHT

def main():
    args = parse_args()
    config = Config(args)
    const = Constants()

    # Initialize presentation and set up master layouts
    prs = Presentation()
    prs.slide_width = const.slide_width
    prs.slide_height = const.slide_height
    layout_0 = setup_and_get_master_layout(prs, "Title Slide", const)
    layout_1 = setup_and_get_master_layout(prs, "Title and Content", const)
    layout_2 = setup_and_get_master_layout(prs, "Title Only", const)

    # Slide 1: Title 
    slide = prs.slides.add_slide(layout_0)
    slide.shapes.title.text = config.title
    # slide.placeholders[1].text = config.subtitle
    tf = slide.placeholders[1].text_frame
    add_bullet_points(tf, config.subtitle.split("\n"), const)

    # Slide 2: Preprocessing & Denoising parameters
    slide = prs.slides.add_slide(layout_1)
    slide.shapes.title.text = "Preprocessing & Denoising parameters"
    tf = slide.placeholders[1].text_frame
    denoising_params = get_denoising_params(config.mat_file)
    lines = [
        f"Polynomial detrending order: {denoising_params['POLY_ORD']}",
        f"Band-pass filter: [{denoising_params['BP_LO']} {denoising_params['BP_HI']}]",
        f"Band-pass after regression (RegBP) v.s. Simultaneous (Simult): {'Simult' if denoising_params['SIMULT'] else 'RegBP'}",
        f"Add quadratic motion parameters: {denoising_params['MOT24']}", 
        f"Number of aCompCor components: {denoising_params['N_ACOMP']}",
        f"Number of ICA-AROMA components: {denoising_params['N_AROMA']}",
        f"Remove the average grey matter (GM) signal: {denoising_params['RM_GMR']}"
    ]
    add_bullet_points(tf, lines, const)

    # Slide 3: QA scores
    slide = prs.slides.add_slide(layout_1)
    slide.shapes.title.text = "QA scores"
    tf = slide.placeholders[1].text_frame
    qa_scores = get_qa_scores(config)
    lines = [
        f"{score_name}: {score_value * 100:.1f}%" if score_value is not None else f"{score_name}: N/A"
        for score_name, score_value in qa_scores.items()
    ]
    add_bullet_points(tf, lines, const)

    # Load summary CSV and loop over each row (contrast) to add slides
    summary_csv = pd.read_csv(os.path.join(config.report_dir, "summary.csv"), header=0)
    summary_csv.sort_values(by=["contrast"], inplace=True)

    for _, row in summary_csv.iterrows():
        contrast   = row["contrast"]
        n_clusters = row["n_clusters"]
        stats_csv  = row["stats_csv"]
        fig_paths  = {
            "Connectome ring": row["ring_jpg"],
            "ROI x ROI Matrix": row["matrix_jpg"],
            "Glass-brain projections": row["brain_jpg"]
        }

        if os.path.exists(stats_csv):
            stats_df = pd.read_csv(stats_csv, header=0)

            # Add a slide for the stats table
            slide = prs.slides.add_slide(layout_2)
            tf = slide.shapes.title.text_frame
            add_contrast_title(tf, contrast, None, const)
            add_stats_table(slide, stats_df, const)

        else:
            # print(f"\nNo significant clusters found for contrast '{contrast}'.")
            # print("Only the slide with the ring figure will be added.")
            fig_paths = { k: v for k, v in fig_paths.items() if k == "Connectome ring" }

        for view_name, fig_path in fig_paths.items():
            if not os.path.exists(fig_path):
                print(f"Warning: Image file '{fig_path}' not found. Skipping.")
                continue

            else:
                # Add a slide for the figure of the current view
                slide = prs.slides.add_slide(layout_2)
                tf = slide.shapes.title.text_frame
                add_contrast_title(tf, contrast, view_name, const)
                # footnote = f"preset = {config.preset}, roi_prefix = {args.roi_prefix}, n_clusters = {n_clusters}"
                footnote = f"n_clusters = {n_clusters}"
                add_figure_and_footnote(slide, fig_path, footnote, const)

    # # Add slide number
    # add_slide_number(prs, const)

    # Save presentation
    prs.save(config.report_path)
    print(f"\nDone! Report saved to {config.report_path}\n")

if __name__ == '__main__':
    main()