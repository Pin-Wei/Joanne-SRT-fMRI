#!/usr/bin/env python3
"""
make_r2r_report.py

Compile CONN ROI-to-ROI 2nd-level results (ring + matrix + glass-brain JPGs +
structured stats CSVs) into:
  1. Markdown report (report.md)
  2. PowerPoint deck (report.pptx)

Output dir is reports/<PROJECT>/ROI-networks_<PRESET>/, so PRESET is encoded
in the folder name. Files inside drop the suffix.

PPTX layout follows /home/aclexp/pinwei/Joanne_SRT_fMRI/tmp/templte.pptx:
  - title slide: centered ctrTitle + subTitle
  - methods slide: title + bullet list
  - per-contrast view slide: rich-text title (kind label color-coded, runs
    as subscript, "minus Random" plain+bold) + large centered figure
    Picture box: off=(1.605", 1.235"), size=(10.114" x 5.518")
    (matches templte.pptx slide3 picture placeholder)

3 view slides per contrast: ring / matrix / glass brain.

Methods reference: https://web.conn-toolbox.org/fmri-methods/cluster-level-inferences

Usage:
  python3 make_r2r_report.py [PROJECT] [PRESET]
    PROJECT defaults to no_PM_260424
    PRESET  defaults to FNC (one of FNC | SPC | TFCE)

  R2R_PROJECT=no_PM_260424 R2R_PRESET=TFCE python3 make_r2r_report.py
"""
from __future__ import annotations

import csv
import os
import re
import sys
from pathlib import Path

PROJECT_ROOT  = Path('/home/aclexp/pinwei/Joanne_SRT_fMRI')
REPORT_ROOT   = PROJECT_ROOT / 'reports'
TEMPLATE_PPTX = PROJECT_ROOT / 'tmp' / 'templte.pptx'

PRESET_NAMES = {
    'FNC':  'Functional Network Connectivity (Jafri et al., 2008) - parametric multivariate cluster-level inference',
    'SPC':  'Spatial Pairwise Clustering (Zalesky et al., 2012) - non-parametric cluster-level inference',
    'TFCE': 'Threshold-Free Cluster Enhancement (Smith & Nichols, 2007) - non-parametric cluster-level inference',
}

PRESET_THRESHOLDS = {
    'FNC':  'connection-level p < .05 (uncorrected); cluster-level p < .05 (FDR-corrected, multivariate omnibus test)',
    'SPC':  'connection-level p < .01 (uncorrected); cluster-level p < .05 (FDR-corrected, mass/intensity)',
    'TFCE': 'connection-level p < .05 (TFCE score); cluster-level p < .05 (uncorrected)',
}

# Plain-text labels (used in markdown, summary table, footer)
CONTRAST_LABELS = {
    'str_main':    'Structured (runs 1-8 averaged) minus Random',
    'str_r12_ran': 'Structured (runs 1 & 2) minus Random',
    'str_r34_ran': 'Structured (runs 3 & 4) minus Random',
    'str_r56_ran': 'Structured (runs 5 & 6) minus Random',
    'str_r78_ran': 'Structured (runs 7 & 8) minus Random',
    'swi_main':    'Switch (runs 3-6) minus Random',
    'swi_r34_ran': 'Switch (runs 3 & 4) minus Random',
    'swi_r56_ran': 'Switch (runs 5 & 6) minus Random',
}

# Rich-text title spec, parsed from contrast id
# (kind_text, kind_rgb, run_text_or_None, plain_runs_text)
KIND_COLORS = {
    'str': (0x00, 0x70, 0xC0),  # blue
    'swi': (0x00, 0xB0, 0x50),  # green
}
KIND_LABELS = {'str': 'Structured', 'swi': 'Switch'}

VIEWS = [
    ('ring',          'Connectome ring'),
    ('matrix',        'ROI x ROI matrix'),
    ('sagittal view', 'Glass-brain projections'),
]


def parse_contrast(c: str):
    """Return (kind, kind_label, kind_rgb, runs_subscript_or_inline, is_subscript)."""
    m = re.match(r'^(str|swi)(?:_r(\d)(\d)_ran|_main)$', c)
    if not m:
        return ('?', c, (0,0,0), '', False)
    kind = m.group(1)
    label = KIND_LABELS[kind]
    rgb = KIND_COLORS[kind]
    if m.group(2) is None:  # *_main
        if kind == 'str':
            inline = ' (runs 1-8) '
        else:
            inline = ' (runs 3-6) '
        return (kind, label, rgb, inline, False)
    sub = f'{m.group(2)}, {m.group(3)}'
    return (kind, label, rgb, sub, True)


# ---------------------------------------------------------------------------
def out_base(project: str, preset: str) -> Path:
    return REPORT_ROOT / project / f'ROI-networks_{preset}'


def fig_path(project: str, preset: str, contrast: str, view_label: str) -> Path:
    return out_base(project, preset) / f'{contrast} ({view_label}).jpg'


def read_summary(project: str, preset: str) -> list[dict]:
    path = out_base(project, preset) / 'summary.csv'
    if not path.exists():
        return []
    with path.open() as fh:
        return list(csv.DictReader(fh))


def read_stats(project: str, preset: str, contrast: str) -> list[dict]:
    path = out_base(project, preset) / 'stats' / f'{contrast}.csv'
    if not path.exists():
        return []
    with path.open() as fh:
        return list(csv.DictReader(fh))


def stat_label(row: dict) -> str:
    name, val, df = row.get('stat_name',''), row.get('stat_value',''), row.get('df','')
    if not name:
        return ''
    if df:
        return f'{name}({df}) = {val}'
    return f'{name} = {val}'


def fmt_p(x: str) -> str:
    if not x:
        return ''
    try:
        v = float(x)
    except ValueError:
        return x
    if v < 0.001:
        return f'{v:.2e}'
    return f'{v:.4f}'


# ---------- Markdown ----------
def generate_markdown(project: str, preset: str) -> Path:
    summary = read_summary(project, preset)
    out = out_base(project, preset) / 'report.md'

    L: list[str] = []
    L += [
        f'# ROI-to-ROI second-level results - {project}',
        '',
        f'**Inference preset:** {preset} - {PRESET_NAMES[preset]}',
        '',
        f'**Thresholds:** {PRESET_THRESHOLDS[preset]}',
        '',
        '**ROI set:** CONN `networks.*` atlas (32 network ROIs; Dosenbach 2010 / Yeo 2011 derivatives)',
        '',
        '**ROI ordering:** hierarchical clustering (Calinski-Harabasz auto cutoff, lambda=0.05)',
        '',
        '## Methods',
        '',
        'gPPI first-level ROI-to-ROI connectivity was estimated for each contrast-of-interest. '
        'Second-level analyses used a one-sample t-test across subjects (AllSubjects factor, contrast `[1]`). '
        f'Inference: CONN standard preset {preset}.',
        '',
        'Reference: <https://web.conn-toolbox.org/fmri-methods/cluster-level-inferences>',
        '',
        '## Summary',
        '',
        '| Contrast | Label | # Networks | # Clusters | # Significant |',
        '|---|---|---:|---:|---:|',
    ]
    for r in summary:
        c = r.get('contrast','')
        L.append('| {} | {} | {} | {} | {} |'.format(
            c, CONTRAST_LABELS.get(c,'-'),
            r.get('n_networks','-'), r.get('n_clusters','-'), r.get('n_sig','-')))
    L += ['', '## Per-contrast results', '']

    for r in summary:
        c     = r.get('contrast','')
        label = CONTRAST_LABELS.get(c, c)
        note  = r.get('note','').strip('"')
        L += [f'### `{c}` - {label}', '']
        if note and note != 'ok':
            L += [f'> **{note}**', '']
            continue

        for view_label, view_desc in VIEWS:
            jpg = fig_path(project, preset, c, view_label)
            if jpg.exists():
                L += [f'**{view_desc}**', '', f'![{c} {view_label}]({jpg.name})', '']
            else:
                L += [f'**{view_desc}** _(figure missing)_', '']

        rows = read_stats(project, preset, c)
        sig_rows = [x for x in rows if x.get('sig') == '1']
        if sig_rows:
            L += [f'**Significant items ({len(sig_rows)}):**', '',
                  '| Type | Cluster | ROI 1 | ROI 2 | Stat | p_unc | p_FDR | p_FWE |',
                  '|---|---:|---|---|---|---:|---:|---:|']
            for x in sig_rows[:50]:
                roi1 = x.get('roi1','') or '-'
                roi2 = x.get('roi2','') or '-'
                if x.get('roi1_xyz'): roi1 = f'{roi1} ({x["roi1_xyz"]})'
                if x.get('roi2_xyz'): roi2 = f'{roi2} ({x["roi2_xyz"]})'
                L.append('| {} | {} | {} | {} | {} | {} | {} | {} |'.format(
                    x.get('type',''), x.get('cluster_id','') or '-',
                    roi1, roi2, stat_label(x),
                    fmt_p(x.get('p_unc','')), fmt_p(x.get('p_FDR','')),
                    fmt_p(x.get('p_FWE',''))))
            if len(sig_rows) > 50:
                L.append(f'| _... ({len(sig_rows)-50} more)_ |')
            L.append('')
        else:
            L += ['_No supra-threshold items at this preset._', '']

    out.write_text('\n'.join(L))
    return out


# ---------- PPTX (template-based) ----------
def _delete_all_slides(prs):
    """Remove all existing slides from a Presentation (so we can repopulate
    using the inherited layouts from templte.pptx)."""
    sldIdLst = prs.slides._sldIdLst
    rId_to_drop = []
    for sldId in list(sldIdLst):
        rId = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        rId_to_drop.append(rId)
        sldIdLst.remove(sldId)
    for rId in rId_to_drop:
        prs.part.drop_rel(rId)


def _layout_by_name(prs, name):
    for lay in prs.slide_layouts:
        if lay.name == name:
            return lay
    raise KeyError(f'layout {name!r} not in template')


def _add_slide_number(slide):
    """Inject a slide-number placeholder + <a:fld type='slidenum'/> into the
    given slide at the template master position (10.167", 6.951", 3" x 0.399"),
    right-aligned. python-pptx's add_slide does NOT auto-include layout
    placeholders, so we add the field shape ourselves."""
    from lxml import etree
    nsmap = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    }
    # 10.167" x 914400 = 9296845 EMU; 6.951" * 914400 = 6358082; 3" = 2743200; 0.399" = 364882
    sp_xml = (
        '<p:sp xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
        '      xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
        '<p:nvSpPr>'
        '<p:cNvPr id="999" name="Slide Number Placeholder"/>'
        '<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>'
        '<p:nvPr><p:ph type="sldNum" sz="quarter" idx="12"/></p:nvPr>'
        '</p:nvSpPr>'
        '<p:spPr/>'
        '<p:txBody>'
        '<a:bodyPr/><a:lstStyle/>'
        '<a:p><a:pPr algn="r"/>'
        '<a:fld id="{EC26BE97-C306-4739-ADC3-79B1BB6A12EE}" type="slidenum">'
        '<a:rPr lang="en-US"/><a:t>#</a:t></a:fld>'
        '</a:p>'
        '</p:txBody>'
        '</p:sp>'
    )
    sp_el = etree.fromstring(sp_xml)
    spTree = slide.shapes._spTree
    spTree.append(sp_el)


def _set_run_props(run, *, bold=None, color=None, sub=False, size_pt=None):
    """Apply font properties to a run, optionally subscripting via baseline."""
    from pptx.util import Pt
    if bold is not None:
        run.font.bold = bold
    if color is not None:
        from pptx.dml.color import RGBColor
        run.font.color.rgb = RGBColor(*color)
    if size_pt is not None:
        run.font.size = Pt(size_pt)
    if sub:
        rPr = run._r.get_or_add_rPr()
        rPr.set('baseline', '-25000')


def _set_title_runs(title_shape, contrast: str, view_desc: str):
    """Populate a TITLE placeholder with rich runs.
       Inherits font size from master (44pt for slide1 layout, 36pt for layout[1] title).
       We do NOT set font.size on runs so the master defaults apply, except we
       shrink the trailing view label."""
    kind, kind_label, rgb, runs_text, is_sub = parse_contrast(contrast)
    tf = title_shape.text_frame
    # Clear all paragraphs except first; clear first paragraph's runs
    p = tf.paragraphs[0]
    p.text = ''
    # remove any extra paragraphs python-pptx may have created
    for extra in list(tf.paragraphs[1:]):
        extra._p.getparent().remove(extra._p)

    r = p.add_run(); r.text = kind_label
    _set_run_props(r, bold=True, color=rgb)

    if is_sub:
        r = p.add_run(); r.text = runs_text
        _set_run_props(r, bold=True, color=rgb, sub=True)
        r = p.add_run(); r.text = ' minus '
    else:
        r = p.add_run(); r.text = runs_text
        r = p.add_run(); r.text = 'minus '

    r = p.add_run(); r.text = 'Random'
    _set_run_props(r, bold=True)

    r = p.add_run(); r.text = f'  -  {view_desc}'
    _set_run_props(r, size_pt=24)


def _img_aspect(path: Path):
    try:
        from PIL import Image
        with Image.open(path) as im:
            w, h = im.size
            return w / h if h else None
    except Exception:
        return None


def _fit_image(slide, jpg: Path, *, area_left, area_top, area_w, area_h, Inches):
    aspect = _img_aspect(jpg) or (11.44 / 6.24)
    area_aspect = area_w / area_h
    if aspect >= area_aspect:
        w = area_w; h = w / aspect
    else:
        h = area_h; w = h * aspect
    left = area_left + (area_w - w) / 2
    top  = area_top  + (area_h - h) / 2
    slide.shapes.add_picture(str(jpg), Inches(left), Inches(top),
                             width=Inches(w), height=Inches(h))


def generate_pptx(project: str, preset: str):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
    except ImportError:
        print('[WARN] python-pptx not installed; skipping PPTX.', file=sys.stderr)
        return None

    summary  = read_summary(project, preset)
    pptx_out = out_base(project, preset) / 'report.pptx'

    # Open template as base — inherits its master, layouts, fonts, positions
    if not TEMPLATE_PPTX.exists():
        print(f'[WARN] template not found: {TEMPLATE_PPTX}', file=sys.stderr)
        return None
    prs = Presentation(str(TEMPLATE_PPTX))
    _delete_all_slides(prs)

    # Slide width/height inherited from template (13.333 x 7.5)

    layout_title   = _layout_by_name(prs, 'Title Slide')        # idx 0
    layout_content = _layout_by_name(prs, 'Title and Content')  # idx 1

    # ---- Slide 1: Title Slide (CENTER_TITLE + SUBTITLE) ----
    slide = prs.slides.add_slide(layout_title)
    _add_slide_number(slide)
    slide.shapes.title.text = f'ROI-to-ROI 2nd-level results'
    sub = slide.placeholders[1]  # SUBTITLE
    tf = sub.text_frame
    tf.text = project
    p = tf.add_paragraph(); p.text = ''
    p = tf.add_paragraph(); p.text = f'Inference preset: {preset}'
    p = tf.add_paragraph(); p.text = PRESET_NAMES[preset]

    # ---- Slide 2: Methods (Title + Content, bullet list) ----
    slide = prs.slides.add_slide(layout_content)
    _add_slide_number(slide)
    slide.shapes.title.text = 'Methods'
    body = slide.placeholders[1]  # OBJECT placeholder
    tf = body.text_frame
    tf.word_wrap = True
    bullets = [
        'gPPI first-level ROI-to-ROI connectivity estimated for each contrast-of-interest.',
        'Second-level: one-sample t-test across subjects (AllSubjects factor, contrast [1]).',
        f'Inference: CONN standard preset {preset} - {PRESET_NAMES[preset]}',
        f'Thresholds: {PRESET_THRESHOLDS[preset]}',
        'ROI set: CONN networks.* atlas (32 network ROIs; Dosenbach 2010 / Yeo 2011 derivatives).',
        'ROI ordering: hierarchical clustering on functional-distance matrix (Calinski-Harabasz auto cutoff, lambda=0.05).',
        'Per-contrast slides: 3 views (ring / matrix / glass-brain projections).',
        'Reference: https://web.conn-toolbox.org/fmri-methods/cluster-level-inferences',
    ]
    for i, txt in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        # Use level 0 (default 28pt from master) for top bullets; the master
        # already supplies bullet glyph + indent.
        p.level = 0

    # ---- Per-contrast slides (Title + Picture, mimic template slide3) ----
    # Picture box from templte.pptx slide3 spPr override:
    #   off=(1.605", 1.235"), ext=(10.114" x 5.518")
    PIC_LEFT, PIC_TOP, PIC_W, PIC_H = 1.605, 1.235, 10.114, 5.518

    for r in summary:
        c     = r.get('contrast','')
        note  = r.get('note','').strip('"')
        rows  = read_stats(project, preset, c)
        sig_rows = [x for x in rows if x.get('sig') == '1']
        n_sig    = len(sig_rows)

        for view_label, view_desc in VIEWS:
            slide = prs.slides.add_slide(layout_content)
            _add_slide_number(slide)
            # We're using Title+Content layout — set rich title runs
            _set_title_runs(slide.shapes.title, c, view_desc)

            # Drop the body placeholder (we're putting an image instead)
            body_ph = slide.placeholders[1]
            body_ph._element.getparent().remove(body_ph._element)

            jpg = fig_path(project, preset, c, view_label)
            if note and note != 'ok':
                err = slide.shapes.add_textbox(Inches(0.5), Inches(7.5/2 - 0.3),
                                               Inches(13.333 - 1.0), Inches(0.6))
                ep = err.text_frame.paragraphs[0]
                ep.text = f'[error] {note}'
                ep.font.size = Pt(18); ep.font.color.rgb = RGBColor(0xCC, 0x00, 0x00)
            elif not jpg.exists():
                miss = slide.shapes.add_textbox(Inches(0.5), Inches(7.5/2 - 0.3),
                                                Inches(13.333 - 1.0), Inches(0.6))
                mp = miss.text_frame.paragraphs[0]
                mp.text = f'[figure missing] {jpg.name}'
                mp.font.size = Pt(16); mp.font.italic = True
            else:
                _fit_image(slide, jpg,
                           area_left=PIC_LEFT, area_top=PIC_TOP,
                           area_w=PIC_W, area_h=PIC_H, Inches=Inches)

            # Small footer above slide-number with preset/sig info
            foot = slide.shapes.add_textbox(Inches(0.5), Inches(7.5 - 0.5),
                                            Inches(13.333 - 4.0), Inches(0.35))
            fp = foot.text_frame.paragraphs[0]
            fp.text = (f'preset={preset}   networks={r.get("n_networks","?")}   '
                       f'clusters={r.get("n_clusters","?")}   significant={n_sig}   view={view_label}')
            fp.font.size = Pt(11)
            fp.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    prs.save(pptx_out)
    return pptx_out


# ---------- main ----------
def main():
    project = sys.argv[1] if len(sys.argv) > 1 else os.environ.get('R2R_PROJECT','no_PM_260424')
    preset  = sys.argv[2] if len(sys.argv) > 2 else os.environ.get('R2R_PRESET', 'FNC')
    if preset not in PRESET_NAMES:
        print(f'Unknown preset {preset}. Use one of: {list(PRESET_NAMES)}')
        sys.exit(1)
    base = out_base(project, preset)
    if not base.exists():
        print(f'[error] output dir does not exist: {base}', file=sys.stderr)
        print('Run run_r2r_networks_bookmarks.m first.', file=sys.stderr)
        sys.exit(2)

    md = generate_markdown(project, preset)
    print(f'Markdown -> {md}')
    pptx = generate_pptx(project, preset)
    if pptx:
        print(f'PPTX     -> {pptx}')


if __name__ == '__main__':
    main()
